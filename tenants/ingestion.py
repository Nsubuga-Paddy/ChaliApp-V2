import base64
import hashlib
import logging
import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from django.conf import settings
from django.db import transaction
from django.utils import timezone
from openai import OpenAI

from .models import KnowledgeChunk, KnowledgeDocument, KnowledgeSourceDocument

logger = logging.getLogger(__name__)

TARGET_CHUNK_TOKENS = 450
CHUNK_OVERLAP_TOKENS = 80
PDF_OCR_MIN_TEXT_CHARS = 40
PDF_OCR_RENDER_SCALE = 2.0


@dataclass
class ExtractedSection:
    text: str
    heading: str = ''
    page_number: int | None = None
    slide_number: int | None = None


def infer_file_type(filename: str) -> str:
    ext = Path(filename).suffix.lower().lstrip('.')
    if ext in {'pdf', 'docx', 'pptx', 'txt'}:
        return ext
    return ''


def index_source_document(source: KnowledgeSourceDocument) -> KnowledgeSourceDocument:
    source.status = KnowledgeSourceDocument.Status.PROCESSING
    source.error_message = ''
    source.save(update_fields=['status', 'error_message', 'updated_at'])

    try:
        with source.file.open('rb') as handle:
            raw = handle.read()
        content_hash = hashlib.sha256(raw).hexdigest()
        sections = extract_source_sections(source, raw)
        chunks = build_chunks(sections)
        if not chunks:
            raise ValueError('No extractable text chunks were produced from this document.')
        embeddings = embed_texts([chunk['text'] for chunk in chunks])

        with transaction.atomic():
            source.content_hash = content_hash
            source.status = KnowledgeSourceDocument.Status.INDEXED
            source.indexed_at = timezone.now()
            source.error_message = ''
            source.save(
                update_fields=[
                    'content_hash',
                    'status',
                    'indexed_at',
                    'error_message',
                    'updated_at',
                ],
            )
            KnowledgeChunk.objects.filter(source_document=source).delete()
            KnowledgeChunk.objects.bulk_create(
                [
                    KnowledgeChunk(
                        company=source.company,
                        source_document=source,
                        chunk_index=index,
                        text=chunk['text'],
                        heading=chunk.get('heading', ''),
                        page_number=chunk.get('page_number'),
                        slide_number=chunk.get('slide_number'),
                        token_count=chunk['token_count'],
                        embedding=embeddings[index] if index < len(embeddings) else [],
                        metadata={
                            'source_title': source.title,
                            'file_type': source.file_type,
                            'is_shareable': source.is_shareable,
                            'origin_url': source.origin_url or '',
                        },
                    )
                    for index, chunk in enumerate(chunks)
                ]
            )
    except Exception as exc:
        logger.exception('Knowledge source indexing failed: %s', source.pk)
        source.status = KnowledgeSourceDocument.Status.FAILED
        source.error_message = str(exc)
        source.save(update_fields=['status', 'error_message', 'updated_at'])
    return source


def schedule_index_source_document(source: KnowledgeSourceDocument) -> None:
    """Queue indexing on Celery after the DB transaction commits."""
    KnowledgeSourceDocument.objects.filter(pk=source.pk).update(
        status=KnowledgeSourceDocument.Status.PENDING,
        error_message='',
    )

    source_id = source.pk

    def _run_index() -> None:
        from .tasks import index_knowledge_source_document

        if getattr(settings, 'CELERY_TASK_ALWAYS_EAGER', False):
            index_knowledge_source_document(source_id)
            return
        try:
            index_knowledge_source_document.delay(source_id)
        except Exception:
            logger.exception(
                'Celery broker unavailable; indexing source %s synchronously',
                source_id,
            )
            index_source_document(
                KnowledgeSourceDocument.objects.get(pk=source_id),
            )

    if getattr(settings, 'CELERY_TASK_ALWAYS_EAGER', False):
        _run_index()
        return

    transaction.on_commit(_run_index)


def _openai_client() -> OpenAI:
    timeout = float(getattr(settings, 'OPENAI_HTTP_TIMEOUT_SECONDS', 120))
    return OpenAI(api_key=settings.OPENAI_API_KEY, timeout=timeout)


def index_legacy_document(document: KnowledgeDocument) -> None:
    sections = [
        ExtractedSection(
            text=document.content,
            heading=document.title,
        )
    ]
    chunks = build_chunks(sections)
    embeddings = embed_texts([chunk['text'] for chunk in chunks])
    with transaction.atomic():
        KnowledgeChunk.objects.filter(legacy_document=document).delete()
        if not document.is_published:
            return
        KnowledgeChunk.objects.bulk_create(
            [
                KnowledgeChunk(
                    company=document.company,
                    legacy_document=document,
                    chunk_index=index,
                    text=chunk['text'],
                    heading=chunk.get('heading', document.title),
                    token_count=chunk['token_count'],
                    embedding=embeddings[index] if index < len(embeddings) else [],
                    metadata={
                        'source_title': document.title,
                        'category': document.category,
                        'tags': document.tag_list,
                        'legacy_document_id': document.id,
                    },
                )
                for index, chunk in enumerate(chunks)
            ]
        )


def extract_source_sections(
    source: KnowledgeSourceDocument,
    raw: bytes,
) -> list[ExtractedSection]:
    if source.file_type == KnowledgeSourceDocument.FileType.TXT:
        return [ExtractedSection(text=raw.decode('utf-8', errors='ignore'), heading=source.title)]
    if source.file_type == KnowledgeSourceDocument.FileType.PDF:
        return extract_pdf_sections(raw)
    if source.file_type == KnowledgeSourceDocument.FileType.DOCX:
        return extract_docx_sections(raw)
    if source.file_type == KnowledgeSourceDocument.FileType.PPTX:
        return extract_pptx_sections(raw)
    raise ValueError(f'Unsupported knowledge source file type: {source.file_type}')


def extract_pdf_sections(raw: bytes) -> list[ExtractedSection]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError('PDF extraction requires pypdf. Install project requirements.') from exc

    reader = PdfReader(BytesIO(raw))
    sections = []
    page_count = len(reader.pages)
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ''
        normalized = normalize_text(text)
        if len(normalized) < PDF_OCR_MIN_TEXT_CHARS:
            logger.info('Running OCR for PDF page %s/%s', index, page_count)
            ocr_text = extract_pdf_page_ocr(raw, index)
            if ocr_text:
                text = ocr_text
        if text.strip():
            sections.append(ExtractedSection(text=text, page_number=index))
    return sections


def extract_pdf_page_ocr(raw: bytes, page_number: int) -> str:
    if not getattr(settings, 'OPENAI_API_KEY', ''):
        logger.warning('Skipping OCR for PDF page %s; OPENAI_API_KEY is not configured.', page_number)
        return ''

    try:
        image_bytes = render_pdf_page_png(raw, page_number)
    except Exception:
        logger.exception('Failed to render PDF page %s for OCR.', page_number)
        return ''

    try:
        client = _openai_client()
        image_b64 = base64.b64encode(image_bytes).decode('ascii')
        model = getattr(settings, 'OPENAI_OCR_MODEL', 'gpt-4o-mini')
        response = client.chat.completions.create(
            model=model,
            temperature=0,
            messages=[
                {
                    'role': 'system',
                    'content': (
                        'You extract text from scanned knowledge-base documents. '
                        'Return only text that is visible on the page. Preserve headings, list structure, '
                        'tables as readable text, and the original language. Do not add explanations.'
                    ),
                },
                {
                    'role': 'user',
                    'content': [
                        {
                            'type': 'text',
                            'text': f'Extract all readable text from page {page_number}.',
                        },
                        {
                            'type': 'image_url',
                            'image_url': {
                                'url': f'data:image/png;base64,{image_b64}',
                            },
                        },
                    ],
                },
            ],
        )
        return normalize_text(response.choices[0].message.content or '')
    except Exception:
        logger.exception('OpenAI OCR failed for PDF page %s.', page_number)
        return ''


def render_pdf_page_png(raw: bytes, page_number: int) -> bytes:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError('Scanned PDF OCR requires PyMuPDF. Install project requirements.') from exc

    with fitz.open(stream=raw, filetype='pdf') as document:
        page = document.load_page(page_number - 1)
        matrix = fitz.Matrix(PDF_OCR_RENDER_SCALE, PDF_OCR_RENDER_SCALE)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        return pixmap.tobytes('png')


def extract_docx_sections(raw: bytes) -> list[ExtractedSection]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError('DOCX extraction requires python-docx. Install project requirements.') from exc

    document = Document(BytesIO(raw))
    sections = []
    heading = ''
    buffer = []

    def flush_buffer() -> None:
        nonlocal buffer
        if buffer:
            sections.append(ExtractedSection(text='\n'.join(buffer), heading=heading))
            buffer = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or '').lower()
        if 'heading' in style_name:
            flush_buffer()
            heading = text
        else:
            buffer.append(text)

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [normalize_text(cell.text) for cell in row.cells]
            row_text = ' | '.join(cell for cell in cells if cell)
            if row_text:
                rows.append(row_text)
        if rows:
            flush_buffer()
            sections.append(ExtractedSection(text='\n'.join(rows), heading=heading or 'Table'))

    flush_buffer()
    return sections


def extract_pptx_sections(raw: bytes) -> list[ExtractedSection]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError('PPTX extraction requires python-pptx. Install project requirements.') from exc

    presentation = Presentation(BytesIO(raw))
    sections = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        notes = getattr(slide, 'notes_slide', None)
        if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
            texts.append(notes.notes_text_frame.text.strip())
        if texts:
            sections.append(
                ExtractedSection(
                    text='\n'.join(texts),
                    heading=f'Slide {slide_index}',
                    slide_number=slide_index,
                )
            )
    return sections


def build_chunks(sections: list[ExtractedSection]) -> list[dict]:
    chunks = []
    for section in sections:
        normalized = normalize_text(section.text)
        if not normalized:
            continue
        tokens = normalized.split()
        if len(tokens) <= TARGET_CHUNK_TOKENS:
            chunks.append(chunk_payload(section, normalized, len(tokens)))
            continue
        step = max(TARGET_CHUNK_TOKENS - CHUNK_OVERLAP_TOKENS, 1)
        for start in range(0, len(tokens), step):
            window = tokens[start : start + TARGET_CHUNK_TOKENS]
            if not window:
                continue
            chunks.append(chunk_payload(section, ' '.join(window), len(window)))
            if start + TARGET_CHUNK_TOKENS >= len(tokens):
                break
    return chunks


def chunk_payload(section: ExtractedSection, text: str, token_count: int) -> dict:
    return {
        'text': text,
        'heading': section.heading,
        'page_number': section.page_number,
        'slide_number': section.slide_number,
        'token_count': token_count,
    }


def normalize_text(text: str) -> str:
    return re.sub(r'\s+', ' ', text or '').strip()


def embed_texts(texts: list[str]) -> list[list[float]]:
    if not texts or not getattr(settings, 'OPENAI_API_KEY', ''):
        return [[] for _ in texts]

    model = getattr(settings, 'OPENAI_EMBEDDING_MODEL', 'text-embedding-3-small')
    try:
        client = _openai_client()
        response = client.embeddings.create(model=model, input=texts)
        by_index = sorted(response.data, key=lambda item: item.index)
        return [list(item.embedding) for item in by_index]
    except Exception:
        logger.exception('Embedding generation failed; falling back to lexical retrieval.')
        return [[] for _ in texts]
